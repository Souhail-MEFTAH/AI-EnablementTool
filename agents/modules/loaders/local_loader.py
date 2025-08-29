from pptx import Presentation
from docx import Document as DocxDocument
from PyPDF2 import PdfReader
import pandas as pd
import json

def load_txt(path):
    with open(path, 'r', encoding='utf-8') as f:
        return f.read()

def load_md(path):
    return load_txt(path)

def load_csv(path):
    df = pd.read_csv(path)
    return df.to_string()

def load_json(path):
    with open(path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    return json.dumps(data, indent=2)

def load_pptx(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def load_docx(path):
    doc = DocxDocument(path)
    return "\n".join([p.text for p in doc.paragraphs])

def load_pdf(path):
    reader = PdfReader(path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def load_xlsx(path):
    df = pd.read_excel(path)
    return df.to_string()

def load_video(path):
    import whisper
    model = whisper.load_model("base")
    result = model.transcribe(str(path))
    return result["text"]

loader_map = {
    ".txt": load_txt,
    ".md": load_md,
    ".csv": load_csv,
    ".json": load_json,
    ".pptx": load_pptx,
    ".docx": load_docx,
    ".pdf": load_pdf,
    ".xlsx": load_xlsx,
    ".mp4": load_video,
    ".mkv": load_video,
    ".avi": load_video,
}

def load_supported_file(path):
    ext = path.suffix.lower()
    loader = loader_map.get(ext)
    if not loader:
        raise ValueError(f"Unsupported file type: {ext}")
    return loader(path)